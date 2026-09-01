#!/usr/bin/env python3
"""Generate a brand partnership pitch deck (.pptx) from JSON.

Input: /workspace/in/brand_deck.json (or first argv) with:
  artist_name, tagline, audience (bullets), stats (label/value),
  value_props, activations, contact
Output: /workspace/out/BrandDeck_<slug>.pptx

Run only via the execute_python sandbox. python-pptx preinstalled.
"""
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt


def load_input() -> dict:
    path = sys.argv[1] if len(sys.argv) > 1 else "/workspace/in/brand_deck.json"
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def add_text_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)


def build(deck: dict) -> None:
    prs = Presentation()
    title = deck.get("artist_name", "Artist")
    slug = deck.get("slug", title.lower().replace(" ", "-").replace("/", "-"))

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = deck.get("tagline", "Partnership Proposal")

    # Audience
    add_text_slide(prs, "Audience", deck.get("audience", []))

    # Stats
    add_text_slide(prs, "Why us (by the numbers)", [
        f"{s.get('label', '')}: {s.get('value', '')}" for s in deck.get("stats", [])
    ] or ["(add stats from Chatmu MCP)"])

    # Value props
    add_text_slide(prs, "What a brand gets", deck.get("value_props", []))

    # Activations
    add_text_slide(prs, "Activation ideas", deck.get("activations", []))

    # Contact
    add_text_slide(prs, "Contact", [
        f"{c.get('role', '')}: {c.get('name', '')} — {c.get('email', '')}"
        for c in deck.get("contacts", [])
    ] or ["(add contact)"])

    out_dir = "/workspace/out"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"BrandDeck_{slug}.pptx")
    prs.save(out_path)
    print(f"Saved {out_path}")


if __name__ == "__main__":
    build(load_input())
